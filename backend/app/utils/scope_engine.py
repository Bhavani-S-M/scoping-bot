# app/utils/scope_engine.py
from __future__ import annotations
import asyncio
import json, re, logging, math, os, tempfile,anyio,pytesseract, openpyxl,tiktoken, pytz, graphviz
from app import models
from calendar import monthrange
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
from io import BytesIO
from PIL import Image
from azure.search.documents.models import VectorizedQuery
from typing import Dict, Any, List
from datetime import datetime, timedelta
from app.utils import azure_blob
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.utils.ai_clients import (
    get_azure_openai_client,
    get_azure_openai_deployment,
    get_azure_openai_embedding_deployment,
    get_search_client,
)

logger = logging.getLogger(__name__)

# Init Azure services
client = get_azure_openai_client()
deployment = get_azure_openai_deployment()
emb_model = get_azure_openai_embedding_deployment()
search_client = get_search_client()

PROJECTS_BASE = "projects"


# Default Role Rates (USD/month)
ROLE_RATE_MAP: Dict[str, float] = {
    "Backend Developer": 3000.0,
    "Frontend Developer": 2800.0,
    "QA Analyst": 1800.0,
    "QA Engineer": 2000.0,
    "Data Engineer": 2800.0,
    "Data Analyst": 2200.0,
    "Data Architect": 3500.0,
    "UX Designer": 2500.0,
    "UI/UX Designer": 2600.0,
    "Project Manager": 3500.0,
    "Cloud Engineer": 3000.0,
    "BI Developer": 2700.0,
    "DevOps Engineer": 3200.0,
    "Security Administrator": 3000.0,
    "System Administrator": 2800.0,
    "Solution Architect": 4000.0,
}

#  helpers
def _strip_code_fences(s: str) -> str:
    m = re.search(r"```(?:json)?(.*?)```", s, flags=re.DOTALL | re.IGNORECASE)
    return m.group(1) if m else s

def _extract_json(s: str) -> dict:
    raw = _strip_code_fences(s or "")
    try:
        return json.loads(raw.strip())
    except Exception:
        start, end = raw.find("{"), raw.rfind("}")
        if start >= 0 and end > start:
            try:
                return json.loads(raw[start:end+1])
            except Exception:
                return {}
        return {}
async def get_rate_map_for_project(db: AsyncSession, project) -> Dict[str, float]:
    """
    Fetch rate cards for the given project/company.
    Falls back to Sigmoid default rates if none exist
    """
    try:
        # If project has company_id, try fetching company-specific rate cards
        if getattr(project, "company_id", None):
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == project.company_id)
            )
            ratecards = result.scalars().all()
            if ratecards:
                return {r.role_name: float(r.monthly_rate) for r in ratecards}

        sigmoid_result = await db.execute(
            select(models.Company).filter(models.Company.name == "Sigmoid")
        )
        sigmoid = sigmoid_result.scalars().first()
        if sigmoid:
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == sigmoid.id)
            )
            sigmoid_rates = result.scalars().all()
            if sigmoid_rates:
                return {r.role_name: float(r.monthly_rate) for r in sigmoid_rates}

    except Exception as e:
        logger.warning(f"Failed to fetch rate cards: {e}")
    return ROLE_RATE_MAP


async def _extract_text_from_files(files: List[dict]) -> str:
    results: List[str] = []

    async def _extract_single(f: dict) -> None:
        try:
            blob_bytes = await azure_blob.download_bytes(f["file_path"])
            suffix = os.path.splitext(f["file_name"])[-1].lower()

            def process_file() -> str:
                content = ""
                try:
                    if suffix == ".pdf":
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            content = extract_pdf_text(tmp_path)
                        finally:
                            os.remove(tmp_path)

                    elif suffix == ".docx":
                        doc = Document(BytesIO(blob_bytes))
                        content = "\n".join(p.text for p in doc.paragraphs)

                    elif suffix == ".pptx":
                        prs = Presentation(BytesIO(blob_bytes))
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    texts.append(shape.text)
                        content = "\n".join(texts)

                    elif suffix in [".xlsx", ".xlsm"]:
                        wb = openpyxl.load_workbook(BytesIO(blob_bytes))
                        sheet = wb.active
                        content = "\n".join(
                            " ".join(str(cell) if cell else "" for cell in row)
                            for row in sheet.iter_rows(values_only=True)
                        )

                    elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
                        img = Image.open(BytesIO(blob_bytes))
                        content = pytesseract.image_to_string(img)

                    else:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as fh:
                                content = fh.read()
                        finally:
                            os.remove(tmp_path)

                except Exception as e:
                    logger.warning(f"Extraction failed for {f['file_name']}: {e}")

                return content.strip()

            text = await anyio.to_thread.run_sync(process_file)

            if text:
                results.append(text)
            else:
                logger.warning(f"Extracted no text from {f['file_name']}")

        except Exception as e:
            logger.warning(f"Failed to extract {f.get('file_name')} (path={f.get('file_path')}): {e}")

    async with anyio.create_task_group() as tg:
        for f in files:
            tg.start_soon(_extract_single, f)

    return "\n\n".join(results)




def _rag_retrieve(query: str, k: int = 3, expand_neighbors: bool = True) -> List[Dict]:
    """
    Retrieve semantically similar chunks for RAG.
    Trims results so total tokens <= model context limit (with buffer).
    """
    global emb_model

    if not search_client or not client:
        return []

    try:
        # Embedding for query
        q_emb = client.embeddings.create(model=emb_model, input=query).data[0].embedding

        # Vector query
        vector_query = VectorizedQuery(
            vector=q_emb,
            fields="text_vector",
        )

        results = search_client.search(
            search_text=query,
            vector_queries=[vector_query],
            select=["chunk_id", "parent_id", "chunk", "title"],
            top=k
        )

        hits = []
        for doc in results:
            hits.append({
                "id": doc["chunk_id"],
                "parent_id": doc.get("parent_id"),
                "content": doc["chunk"],
                "title": doc.get("title")
            })

        # Expand neighbors
        if expand_neighbors and hits:
            expanded = []
            for h in hits:
                expanded.append(h)
                if "_" in h["id"]:
                    try:
                        base, idx = h["id"].rsplit("_", 1)
                        idx = int(idx)
                        for nid in [f"{base}_{idx-1}", f"{base}_{idx+1}"]:
                            try:
                                neighbor = search_client.get_document(nid)
                                if neighbor:
                                    expanded.append({
                                        "id": neighbor["chunk_id"],
                                        "parent_id": neighbor.get("parent_id"),
                                        "content": neighbor["chunk"],
                                        "title": neighbor.get("title")
                                    })
                            except Exception:
                                pass
                    except Exception:
                        pass
            hits = expanded

        # Group by parent_id
        grouped = {}
        for h in hits:
            grouped.setdefault(h["parent_id"], []).append(h)

        # ---- Token budget check ----
        model_name = deployment
        tokenizer = tiktoken.encoding_for_model(model_name)
        context_limit = 128000
        max_tokens = context_limit - 4000 
        used_tokens = 0

        ordered = []
        for pid, docs in grouped.items():
            safe_chunks = []
            for d in docs:
                tokens = len(tokenizer.encode(d["content"]))
                if used_tokens + tokens > max_tokens:
                    break
                safe_chunks.append(d)
                used_tokens += tokens

            if safe_chunks:
                ordered.append({
                    "parent_id": pid,
                    "chunks": safe_chunks,
                    "title": safe_chunks[0].get("title")
                })

        logger.info(f"RAG retrieve kept {used_tokens} tokens (limit {max_tokens})")

        return ordered

    except Exception as e:
        logger.warning(f"RAG retrieve failed: {e}")
        return []


# ---------- Prompt ----------
def _build_scope_prompt(rfp_text: str, kb_chunks: List[str], project=None, model_name: str = "gpt-4o", questions_context: str | None = None) -> str:
    import datetime, tiktoken

    # Tokenizer
    tokenizer = tiktoken.encoding_for_model(model_name)

    # Safe token budget (GPT-4o = 128k, keep ~4k for completion & system messages)
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # Trim RFP text
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 3000:   # still enforce cap on huge RFPs
        rfp_tokens = rfp_tokens[:3000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # Trim KB context
    safe_kb_chunks = []
    for ch in kb_chunks or []:
        tokens = tokenizer.encode(ch)
        if used_tokens + len(tokens) > max_total_tokens:
            break
        safe_kb_chunks.append(ch)
        used_tokens += len(tokens)

    kb_context = "\n\n".join(safe_kb_chunks) if safe_kb_chunks else "(no KB context found)"

    # ---------- Project user fields ----------
    name = (getattr(project, "name", "") or "").strip()
    domain = (getattr(project, "domain", "") or "").strip()
    complexity = (getattr(project, "complexity", "") or "").strip()
    tech_stack = (getattr(project, "tech_stack", "") or "").strip()
    use_cases = (getattr(project, "use_cases", "") or "").strip()
    compliance = (getattr(project, "compliance", "") or "").strip()
    duration = str(getattr(project, "duration", "") or "").strip()

    user_context = (
        "Some overview fields have been provided by the user.\n"
        "Treat these user-provided values as the source of truth.\n"
        "Only fill in fields that are blank — do NOT overwrite the given values.\n\n"
        f"Project Name: {name or '(infer if missing)'}\n"
        f"Domain: {domain or '(infer if missing)'}\n"
        f"Complexity: {complexity or '(infer if missing)'}\n"
        f"Tech Stack: {tech_stack or '(infer if missing)'}\n"
        f"Use Cases: {use_cases or '(infer if missing)'}\n"
        f"Compliance: {compliance or '(infer if missing)'}\n"
        f"Duration (months): {duration or '(infer if missing)'}\n\n"
    )

    today_str = datetime.date.today().isoformat()

    # ---------- Final Prompt ----------
    return (
        "You are an expert AI project planner.\n"
        "Use the RFP/project text as the **primary source**, but enrich missing fields "
        "with the Knowledge Base context (if relevant).\n"
        "Return ONLY valid JSON (no prose, no markdown, no commentary).\n\n"
        "Output schema:\n"
        "{\n"
        '  "overview": {\n'
        '    "Project Name": string,\n'
        '    "Domain": string,\n'
        '    "Complexity": string,\n'
        '    "Tech Stack": string,\n'
        '    "Use Cases": string,\n'
        '    "Compliance": string,\n'
        '    "Duration": number\n'
        "  },\n"
        '  "activities": [\n'
        '    {\n'
        '      "ID": int,\n'
        '      "Activities": string,\n'
        '      "Description": string | null,\n'
        '      "Owner": string | null,\n'
        '      "Resources": string | null,\n'
        '      "Start Date": "yyyy-mm-dd",\n'
        '      "End Date": "yyyy-mm-dd",\n'
        '      "Effort Months": int\n'
        "    }\n"
        "  ],\n"
        '  "resourcing_plan": []\n'
        "}\n\n"
        f"-Rules:\n"
        f"- The first activity must always start today ({today_str}).\n"
        "- **Allow maximum parallel execution**: activities with no dependency must overlap instead of running sequentially."
        "- Project duration must always be **under 12 months**.\n"
        "- Auto-calculate **End Date = Start Date + Effort Months**.\n"
        "- Auto-calculate **overview.Duration** as the total span in months from the earliest Start Date to the latest End Date.\n"
        "- `Complexity` should be simple, medium, or large based on duration of project.\n"
        "- **Always assign at least one Resource**."
        "- Distinguish `Owner` (responsible lead role) and `Resources` (supporting roles)."
        "- `Owner` and `Resources` must be valid IT roles (e.g., Backend Developer, AI Engineer, QA Engineer, etc.)."
        "- `Owner` is always a role who manages that particular activity (not a personal name).\n"
        "- `Resources` must contain only roles which are required for that particular activity, distinct from `Owner`.\n"
        "- If `Resources` is missing, fallback to the same `Owner` role.\n"
        "- Use less resources as much as possible.\n"
        "- Effort Months should be small integers 0.5 to 1.5 months not more than this.\n"
        "- IDs must start from 1 and increment sequentially.\n"
        "- If the RFP or Knowledge Base text lacks detail, infer the missing pieces logically."
        "- Include all relevant roles and activities that ensure delivery of the project scope."
        "- Keep all field names exactly as in the schema.\n"

        f"{user_context}"
        f"RFP / Project Files Content:\n{rfp_text}\n\n"
        f"Knowledge Base Context (for enrichment only):\n{kb_context}\n"
        f"Clarification Q&A (User-confirmed answers take highest priority)\n"
        f"Use these answers to override or clarify any ambiguous or conflicting information.\n"
        f"Do NOT hallucinate beyond these facts.\n\n"
        f"{questions_context}\n"
    )
def _build_questionnaire_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    """
    Build a focused prompt to generate categorized questions
    grouped under relevant topics (Architecture, Data, Security, etc.)
    """
    name = getattr(project, "name", "Unnamed Project")
    domain = getattr(project, "domain", "General")
    tech = getattr(project, "tech_stack", "Modern Web Stack")
    compliance = getattr(project, "compliance", "General")
    duration = getattr(project, "duration", "TBD")

    return f"""
You are an expert business analyst who reviews RFPs and creates structured questionnaires
to clarify requirements before project scoping.

Based on the following information, generate **categorized questions** grouped logically.

Project Context:
- Project Name: {name}
- Domain: {domain}
- Tech Stack: {tech}
- Compliance: {compliance}
- Duration: {duration}

RFP Content:
{rfp_text}

Knowledge Base Context:
{kb_chunks}

---

Return ONLY valid JSON in the following format:

{{
  "questions": [
    {{
      "category": "Architecture",
      "items": [
        {{
          "question": "What is the preferred deployment model?",
          "user_understanding": "",
          "comment": ""
        }},
        {{
          "question": "Do you need auto-scaling or load balancing?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }},
    {{
      "category": "Data & Security",
      "items": [
        {{
          "question": "Will sensitive data be stored or processed?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }}
  ]
}}

Rules:
- Group questions by meaningful categories (Architecture, Data, Integration, Compliance, Delivery, etc.)
- Each category must have at least 2 questions.
- Every question must be clear, specific, and require a short textual answer.
- Always include empty strings for 'user_understanding' and 'comment'.
- Return ONLY valid JSON (no markdown, no explanations).
"""
def _extract_questions_from_text(raw_text: str) -> list[dict]:
    """
    Extract categorized questions from JSON or text.
    Output format:
    [
      {
        "category": "Architecture",
        "items": [
          {"question": "...", "user_understanding": "", "comment": ""}
        ]
      }
    ]
    """
    try:
        parsed = _extract_json(raw_text)

        # Case 1: Proper JSON with nested categories
        if isinstance(parsed, dict) and "questions" in parsed:
            qdata = parsed["questions"]
            if isinstance(qdata, list) and all(isinstance(x, dict) for x in qdata):
                # check if already nested structure
                if "items" in qdata[0]:
                    normalized = []
                    for cat in qdata:
                        normalized.append({
                            "category": cat.get("category", "General"),
                            "items": [
                                {
                                    "question": i.get("question", ""),
                                    "user_understanding": i.get("user_understanding", ""),
                                    "comment": i.get("comment", "")
                                } for i in cat.get("items", [])
                            ]
                        })
                    return normalized

                # Otherwise, flat → group by category
                grouped = {}
                for q in qdata:
                    cat = q.get("category", "General") if isinstance(q, dict) else "General"
                    que = q.get("question", q) if isinstance(q, dict) else str(q)
                    grouped.setdefault(cat, []).append({
                        "question": que,
                        "user_understanding": "",
                        "comment": ""
                    })
                return [{"category": c, "items": lst} for c, lst in grouped.items()]

        # Case 2: List of plain questions
        if isinstance(parsed, list):
            return [{
                "category": "General",
                "items": [{"question": str(q), "user_understanding": "", "comment": ""} for q in parsed]
            }]
    except Exception:
        pass

    # Fallback — parse raw text
    current_cat = "General"
    grouped: dict[str, list] = {}
    for line in raw_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if re.match(r"^(#+\s*)?([A-Z][A-Za-z\s&/]+):?$", line) and not line.endswith("?"):
            current_cat = re.sub(r"^#+\s*", "", line).strip(": ").strip()
            continue
        if "?" in line:
            qtext = re.sub(r"^\d+[\).\s]+", "", line).strip()
            grouped.setdefault(current_cat, []).append({
                "question": qtext,
                "user_understanding": "",
                "comment": ""
            })

    return [{"category": c, "items": lst} for c, lst in grouped.items()]
async def generate_project_questions(db: AsyncSession, project) -> dict:
    """
    Generate a categorized questionnaire for the given project using Azure OpenAI.
    Saves the questions.json file in Azure Blob.
    """
    if client is None or deployment is None:
        logger.warning("Azure OpenAI not configured — skipping question generation")
        return {"questions": []}

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        if getattr(project, "files", None):
            files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            if files:
                rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"Failed to extract RFP for questions: {e}")

    # ---------- Retrieve Knowledge Base ----------
    kb_results = _rag_retrieve(rfp_text or project.name or project.domain)
    kb_chunks = [ch["content"] for group in kb_results for ch in group["chunks"]] if kb_results else []

    # ---------- Build prompt ----------
    prompt = _build_questionnaire_prompt(rfp_text, kb_chunks, project)

    # ---------- Query Azure OpenAI ----------
    try:
        resp = await anyio.to_thread.run_sync(
            lambda: client.chat.completions.create(
                model=deployment,
                messages=[
                    {"role": "system", "content": "You are an expert business analyst."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.4,
            )
        )

        raw_text = resp.choices[0].message.content.strip()
        questions = _extract_questions_from_text(raw_text)
        total_q = sum(len(cat["items"]) for cat in questions)
        logger.info(f" Generated {total_q} questions under {len(questions)} categories for project {project.id}")

        # ---------- Save to Blob Storage ----------
        blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        try:
            await azure_blob.upload_bytes(
                json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
            )

            db_file = models.ProjectFile(
                project_id=project.id,
                file_name="questions.json",
                file_path=blob_name,
            )

            db.add(db_file)
            await db.commit()
            await db.refresh(db_file)

            logger.info(f" Saved questions.json for project {project.id}")
        except Exception as e:
            logger.warning(f"Failed to upload questions.json: {e}")

        return {"questions": questions}

    except Exception as e:
        logger.error(f" Question generation failed: {e}")
        return {"questions": []}
    
# Update questions.json with user input answers
async def update_questions_with_user_input(
    db: AsyncSession, project, user_answers: dict
) -> dict:
    """
    Merge user answers into the existing questions.json for the given project.
    Example user_answers structure:
    {
      "Architecture": {
         "What is the preferred deployment model?": "Cloud-based",
         "Do you need auto-scaling or load balancing?": "Yes, via AKS"
      },
      "Data & Security": {
         "Will sensitive data be stored or processed?": "Yes, PII data"
      }
    }
    """
    from app.utils import azure_blob

    blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
    try:
        # Load current questions.json
        q_bytes = await azure_blob.download_bytes(blob_name)
        q_json = json.loads(q_bytes.decode("utf-8"))
        questions = q_json.get("questions", [])

        # Merge answers into the structure
        for cat in questions:
            cat_name = cat.get("category")
            for item in cat.get("items", []):
                q_text = item.get("question")
                ans = (
                    user_answers.get(cat_name, {}).get(q_text)
                    if user_answers.get(cat_name)
                    else None
                )
                if ans:
                    item["user_understanding"] = ans

        # Upload updated JSON to Blob
        new_bytes = json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8")
        await azure_blob.upload_bytes(new_bytes, blob_name)
        logger.info(f" Updated questions.json with user input for project {project.id}")

        #  Save / update DB record
        db_file = models.ProjectFile(
            project_id=project.id,
            file_name="questions.json",
            file_path=blob_name,
        )
        db.add(db_file)
        await db.commit()
        await db.refresh(db_file)

        return {"questions": questions}

    except Exception as e:
        logger.error(f"Failed to update questions.json with user input: {e}")
        return {}




def _parse_date_safe(val: Any, fallback: datetime = None) -> datetime:
    """Try to parse a date string; return fallback if invalid."""
    if not val:
        return fallback
    try:
        return datetime.strptime(str(val), "%Y-%m-%d")
    except Exception:
        return fallback

def _safe_str(val: Any) -> str:
    return str(val).strip() if val is not None else ""

    
def _build_architecture_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    name = (getattr(project, "name", "") or "Untitled Project").strip()
    domain = (getattr(project, "domain", "") or "General").strip()
    tech = (getattr(project, "tech_stack", "") or "Modern Web + Cloud Stack").strip()

    return f"""
    You are a **senior enterprise solution architect** tasked with designing a *tailored cloud system architecture diagram*
    strictly based on the provided RFP and contextual knowledge.

    ### PROJECT CONTEXT
    - **Project Name:** {name}
    - **Domain:** {domain}
    - **Tech Stack:** {tech}

    ### RFP SUMMARY
    {rfp_text}

    ### KNOWLEDGE BASE CONTEXT
    {kb_chunks}

    ---

    ###  STEP 1 — Reasoning (Internal)
    Analyze the provided RFP and knowledge base to:
    1. Identify all domain-specific **entities, systems, or technologies** mentioned or implied.
    2. Categorize each component into the most appropriate architecture layer:
    - Frontend (UI/Apps)
    - Backend (Services/APIs)
    - Data (Databases, Storage, External APIs)
    - AI/Analytics (ML, Insights, NLP, Recommendations)
    - Security/Monitoring/DevOps (IAM, Key Vault, CI/CD, Logging)
    3. Infer **connections and data flows** between components (e.g., API requests, pipelines, message queues).
    4. Skip any layers not relevant to this RFP.

    You will use this reasoning to build the architecture — but **do not include this reasoning** in your final output.

    ---

    ###  STEP 2 — Graphviz DOT Output
    Generate **only valid Graphviz DOT code** representing the inferred architecture.

    Follow these rules strictly:
    - Begin with: `digraph Architecture {{`
    - End with: `}}`
    - Use **horizontal layout** → `rankdir=LR`
    - Include **only relevant clusters** (omit unused layers)
    - Keep ≤ 15 nodes total
    - Use **orthogonal edges** (`splines=ortho`)
    - Each node label must clearly represent an actual system, service, or tool
    - Logical flow should follow Frontend → Backend → Data → AI → Security (only if applicable)
    -  **Ensure data layers both receive and provide information** — show arrows *into* and *out of* data/storage nodes if analytics, AI, or reporting components exist.

    ---

    ### VISUAL STYLE
    - **Graph:** dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.3
    - **Clusters:** style="filled,rounded", fontname="Helvetica-Bold", fontsize=13
    - **Node Shapes and Colors:**
    - Frontend → `box`, pastel blue (`fillcolor="#E3F2FD"`)
    - Backend/API → `box3d`, pastel green (`fillcolor="#E8F5E9"`)
    - Data/Storage → `cylinder`, pastel yellow (`fillcolor="#FFFDE7"`)
    - AI/Analytics → `ellipse`, pastel purple (`fillcolor="#F3E5F5"`)
    - Security/Monitoring → `diamond`, gray (`fillcolor="#ECEFF1"`)
    - **Edges:** color="#607D8B", penwidth=1.5, arrowsize=0.9

    ---

    ###  STEP 3 — Domain Intelligence (Auto-Enrichment)
    If applicable, automatically enrich the architecture using these domain patterns:

    - **FinTech** → Payment Gateway, Fraud Detection, KYC/AML Service, Ledger DB
    - **HealthTech** → Patient Portal, EHR System, FHIR API, HIPAA Compliance Layer
    - **GovTech** → Citizen Portal, Secure API Gateway, Compliance & Audit Logging
    - **AI/ML Projects** → Model API, Embedding Store, Training Pipeline, Monitoring Service
    - **Data Platforms** → ETL Pipeline, Data Lake, BI Dashboard
    - **Enterprise SaaS** → Tenant Manager, Auth Service, Billing & Subscription Module

    Include these elements **only if they logically fit** the RFP description.

    ---

    ###  STEP 4 — OUTPUT RULES
    - Output *only* the Graphviz DOT syntax — **no markdown**, **no reasoning**, **no commentary**
    - The final response should be a single valid DOT diagram ready for rendering
    """

async def _generate_fallback_architecture(
    db: AsyncSession,
    project,
    blob_base_path: str
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate and upload a default fallback architecture diagram (4-layer generic layout).
    Triggered when Azure OpenAI or Graphviz generation fails.
    """
    logger.warning(" Using fallback default architecture layout")

    # --- Default DOT diagram ---
    fallback_dot = """
digraph Architecture {
    rankdir=LR;
    graph [dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.2, splines=ortho];
    node [style="rounded,filled", fontname="Helvetica-Bold", fontsize=13, penwidth=1.2];

    subgraph cluster_frontend {
        label="Frontend / User Touchpoints";
        style="filled,rounded"; fillcolor="#E3F2FD";
        web[label="Web App (React / Angular)", shape=box, fillcolor="#BBDEFB"];
        mobile[label="Mobile App", shape=box, fillcolor="#BBDEFB"];
    }

    subgraph cluster_backend {
        label="Backend / Services";
        style="filled,rounded"; fillcolor="#E8F5E9";
        api[label="Core API (FastAPI / Node.js)", shape=box3d, fillcolor="#C8E6C9"];
        auth[label="Auth Service", shape=box3d, fillcolor="#C8E6C9"];
    }

    subgraph cluster_data {
        label="Data / Storage";
        style="filled,rounded"; fillcolor="#FFFDE7";
        db[label="Database (PostgreSQL)", shape=cylinder, fillcolor="#FFF9C4"];
        blob[label="Blob Storage", shape=cylinder, fillcolor="#FFF9C4"];
    }

    subgraph cluster_ai {
        label="AI / Analytics";
        style="filled,rounded"; fillcolor="#F3E5F5";
        ai[label="AI Engine / Insights", shape=ellipse, fillcolor="#E1BEE7"];
        dashboard[label="BI Dashboard", shape=ellipse, fillcolor="#E1BEE7"];
    }

    # Data flow
    web -> api -> db;
    mobile -> api;
    db -> ai -> dashboard;
    api -> auth;
}
"""

    # --- Render DOT → PNG & SVG ---
    tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
    try:
        graph = graphviz.Source(fallback_dot, engine="dot")
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Fallback Graphviz rendering failed: {e}")
        return None, ""

    # --- Upload both files to Azure Blob ---
    blob_name_png = f"{blob_base_path}/architecture_fallback_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_fallback_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)
        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # --- Save both records in DB ---
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Fallback architecture diagrams stored for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png



async def generate_architecture(
    db: AsyncSession,
    project,
    rfp_text: str,
    kb_chunks: List[str],
    blob_base_path: str,
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate a visually clean, context-aware architecture diagram (PNG & SVG)
    from RFP + KB context using Azure OpenAI + Graphviz.
    Uses dynamic prompts that adapt layers automatically (no static template).
    Includes retry logic, sanitization, validation, and fallback diagram.
    """
    if client is None or deployment is None:
        logger.warning(" Azure OpenAI not configured — skipping architecture generation")
        return None, ""

    prompt = _build_architecture_prompt(rfp_text, kb_chunks, project)

    # ---------- Step 1: Ask Azure OpenAI for Graphviz DOT code ----------
    async def _generate_dot_from_ai(retry: int = 0) -> str:
        """Call Azure OpenAI with retry logic."""
        try:
            resp = await anyio.to_thread.run_sync(
                lambda: client.chat.completions.create(
                    model=deployment,
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "You are an expert cloud solution architect. "
                                "You must output ONLY valid Graphviz DOT syntax — no markdown or commentary."
                            ),
                        },
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.7,
                )
            )

            ai_output = resp.choices[0].message.content.strip()

            # Log reasoning part (everything before 'digraph')
            if "digraph" in ai_output:
                reasoning = ai_output.split("digraph", 1)[0].strip()
                if reasoning:
                    logger.info(f" Architecture reasoning summary:\n{reasoning[:800]}")
                ai_output = "digraph" + ai_output.split("digraph", 1)[1]

            return ai_output
        except Exception as e:
            if retry < 2:
                logger.warning(f" AI call failed (retry {retry+1}/3): {e}")
                await anyio.sleep(2)
                return await _generate_dot_from_ai(retry + 1)
            logger.error(f" Azure OpenAI architecture generation failed after retries: {e}")
            return ""

    dot_code = await _generate_dot_from_ai()
    if not dot_code:
        logger.warning(" No DOT code returned by AI — generating fallback diagram")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 2: Clean & sanitize DOT ----------
    dot_code = re.sub(r"```[a-zA-Z]*", "", dot_code).replace("```", "").strip()
    dot_code = dot_code.strip("`").strip()
    dot_code = re.sub(r"(?i)^graph\s", "digraph ", dot_code)

    # Fix brace mismatch
    open_braces = dot_code.count("{")
    close_braces = dot_code.count("}")
    if open_braces > close_braces:
        dot_code += "}" * (open_braces - close_braces)
    elif close_braces > open_braces:
        dot_code = "digraph Architecture {\n" + dot_code

    if not dot_code.lower().startswith("digraph"):
        dot_code = f"digraph Architecture {{\n{dot_code}\n}}"

    # Remove control characters
    dot_code = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", "", dot_code)

    # ---------- Step 3: Do NOT override GPT’s style ----------
    # Keep GPT’s own clusters, nodes, and colors — just ensure it's syntactically valid
    # (Old static preamble removed intentionally)

    # ---------- Step 4: Render DOT → PNG & SVG ----------
    try:
        tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
        graph = graphviz.Source(dot_code, engine="dot")

        # Render both PNG and SVG for better clarity
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Graphviz rendering failed: {e}\n--- DOT Snippet ---\n{dot_code[:800]}")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 5: Upload PNG to Azure Blob ----------
    blob_name_png = f"{blob_base_path}/architecture_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)

        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # ---------- Step 6: Replace old architecture file ----------
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "architecture.png",
        )
    )
    old_file = result.scalars().first()
    if old_file:
        try:
            await azure_blob.delete_blob(old_file.file_path)
            await db.delete(old_file)
            await db.commit()
        except Exception as e:
            logger.warning(f" Failed to delete old architecture.png: {e}")

    # ---------- Step 7: Save new ProjectFile records (PNG + SVG) ----------
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Architecture diagrams stored successfully for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png


# --- Cleaner ---
async def clean_scope(db: AsyncSession, data: Dict[str, Any], project=None) -> Dict[str, Any]:
    if not isinstance(data, dict):
        return {}

    ist = pytz.timezone("Asia/Kolkata")
    today = datetime.now(ist).replace(hour=0, minute=0, second=0, microsecond=0)

    activities: List[Dict[str, Any]] = []
    start_dates, end_dates = [], []
    role_month_map: Dict[str, Dict[str, float]] = {}
    role_order: List[str] = []

    # --- Helper: compute monthly allocation based on actual days in month ---
    def month_effort(s: datetime, e: datetime) -> Dict[str, float]:
        cur = s
        month_eff = {}
        while cur <= e:
            year, month = cur.year, cur.month
            days_in_month = monthrange(year, month)[1]
            start_day = cur.day if cur.month == s.month else 1
            end_day = e.day if cur.month == e.month else days_in_month
            days_count = end_day - start_day + 1
            month_eff[f"{cur.strftime('%b %Y')}"] = round(days_count / 30.0, 2)
            # move to next month
            if month == 12:
                cur = datetime(year + 1, 1, 1)
            else:
                cur = datetime(cur.year, cur.month + 1, 1)
        return month_eff

    # --- Process activities ---
    for idx, a in enumerate(data.get("activities") or [], start=1):
        owner = a.get("Owner") or "Unassigned"

        # Parse dependencies
        raw_deps = [d.strip() for d in str(a.get("Resources") or "").split(",") if d.strip()]

        # Remove owner from resources if duplicated
        raw_deps = [r for r in raw_deps if r.lower() != owner.lower()]

        # Owner always included, then other resources
        roles = [owner] + raw_deps

        s = _parse_date_safe(a.get("Start Date"), today)
        e = _parse_date_safe(a.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        # --- allocate per month (no splitting among roles) ---
        month_alloc = month_effort(s, e)
        for role in roles:
            if role not in role_month_map:
                role_month_map[role] = {}
                role_order.append(role)
            for m, eff in month_alloc.items():
                role_month_map[role][m] = role_month_map[role].get(m, 0.0) + eff

        dur_days = max(1, (e - s).days)
        activities.append({
            "ID": idx,
            "Activities": _safe_str(a.get("Activities")),
            "Description": _safe_str(a.get("Description")),
            "Owner": owner,
            "Resources": ", ".join(raw_deps), 
            "Start Date": s,
            "End Date": e,
            "Effort Months": round(dur_days / 30.0, 2),
        })

        start_dates.append(s)
        end_dates.append(e)

        # --- Sort activities ---
    activities.sort(key=lambda x: x["Start Date"])
    for idx, a in enumerate(activities, start=1):
        a["ID"] = idx
        a["Start Date"] = a["Start Date"].strftime("%Y-%m-%d")
        a["End Date"] = a["End Date"].strftime("%Y-%m-%d")

    # --- Project span & month labels (Month 1, Month 2, ...) ---
    min_start = min(start_dates) if start_dates else today
    max_end = max(end_dates) if end_dates else min_start
    duration = max(1.0, round(max(1, (max_end - min_start).days) / 30.0, 2))
    total_months = max(1, math.ceil((max_end - min_start).days / 30.0))

    month_labels = [f"Month {i}" for i in range(1, total_months + 1)]

    # --- Build per-role, per-month day usage ---
    role_month_usage: Dict[str, Dict[str, float]] = {r: {m: 0.0 for m in month_labels} for r in role_order}

    # Compute total active days per relative month window
    for act in activities:
        s = _parse_date_safe(act.get("Start Date"), today)
        e = _parse_date_safe(act.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        involved_roles = [act.get("Owner") or "Unassigned"] + [
            r.strip() for r in str(act.get("Resources") or "").split(",") if r.strip()
        ]

        for m_idx in range(total_months):
            rel_start = min_start + timedelta(days=m_idx * 30)
            rel_end = min_start + timedelta(days=(m_idx + 1) * 30)

            # overlap between activity and this relative month window
            overlap_start = max(s, rel_start)
            overlap_end = min(e, rel_end)
            overlap_days = max(0, (overlap_end - overlap_start).days)

            if overlap_days > 0:
                for r in involved_roles:
                    if r not in role_month_usage:
                        role_month_usage[r] = {ml: 0.0 for ml in month_labels}
                    role_month_usage[r][f"Month {m_idx + 1}"] += overlap_days

    # --- Convert days to effort with 4-tier partial-month logic ---
    for r, months in role_month_usage.items():
        for m, days in months.items():
            if days > 21:
                months[m] = 1.0
            elif 15 <= days <= 21:
                months[m] = 0.75
            elif 8 <= days < 15:
                months[m] = 0.5
            elif 1 <= days < 8:
                months[m] = 0.25
            else:
                months[m] = 0.0

    try:
        if db:
            ROLE_RATE_MAP_DYNAMIC = await get_rate_map_for_project(db, project)
        else:
            ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP
    except Exception as e:
        logger.warning(f"Rate map fallback due to error: {e}")
        ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP


    # --- Build final resourcing plan ---
    resourcing_plan = []
    for idx, role in enumerate(role_order, start=1):
        month_efforts = role_month_usage.get(role, {m: 0 for m in month_labels})
        total_effort = sum(month_efforts.values())
        rate = ROLE_RATE_MAP_DYNAMIC.get(role, ROLE_RATE_MAP.get(role, 2000.0))
        cost = round(total_effort * rate, 2)
        plan_entry = {
            "ID": idx,
            "Resources": role,
            "Rate/month": rate,
            **month_efforts,
            "Efforts": total_effort,
            "Cost": cost,
        }
        resourcing_plan.append(plan_entry)


    # --- Overview ---
    ov = data.get("overview") or {}
    data["overview"] = {
        "Project Name": _safe_str(ov.get("Project Name") or getattr(project, "name", "Untitled Project")),
        "Domain": _safe_str(ov.get("Domain") or getattr(project, "domain", "")),
        "Complexity": _safe_str(ov.get("Complexity") or getattr(project, "complexity", "")),
        "Tech Stack": _safe_str(ov.get("Tech Stack") or getattr(project, "tech_stack", "")),
        "Use Cases": _safe_str(ov.get("Use Cases") or getattr(project, "use_cases", "")),
        "Compliance": _safe_str(ov.get("Compliance") or getattr(project, "compliance", "")),
        "Duration": duration,
        "Generated At": datetime.now(ist).strftime("%Y-%m-%d %H:%M %Z"),
    }
    try:
        if getattr(project, "company", None):
            data["overview"]["Currency"] = getattr(project.company, "currency", "USD")
        else:
            data["overview"]["Currency"] = "USD"
    except Exception:
        data["overview"]["Currency"] = "USD"


    data["activities"] = activities
    data["resourcing_plan"] = resourcing_plan
    return data


async def generate_project_scope(db: AsyncSession, project) -> dict:
    """
    Generate project scope + architecture diagram + store architecture in DB + return combined JSON.
    """

    # ✅ Ensure the project has a valid company reference (fallback to Sigmoid)
    if not getattr(project, "company_id", None):
        from app.utils import ratecards
        sigmoid = await ratecards.get_or_create_sigmoid_company(db)
        project.company_id = sigmoid.id
        await db.commit()
        await db.refresh(project)
        logger.info(f"🔗 Linked project {project.id} to Sigmoid company as fallback")

    if client is None or deployment is None:
        logger.warning("Azure OpenAI not configured")
        return {}

    model_name = deployment
    tokenizer = tiktoken.encoding_for_model(model_name)
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        files: List[dict] = []
        if getattr(project, "files", None):
            try:
                files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            except Exception as e:
                logger.warning(f" Could not access project.files: {e}")
                files = []
        if files:
            rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"File extraction for project {getattr(project, 'id', None)} failed: {e}")

    # ---------- Trim RFP text ----------
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 5000:
        rfp_tokens = rfp_tokens[:5000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # ---------- Retrieve KB context ----------
    fallback_fields = [
        getattr(project, "name", None),
        getattr(project, "domain", None),
        getattr(project, "complexity", None),
        getattr(project, "tech_stack", None),
        getattr(project, "use_cases", None),
        getattr(project, "compliance", None),
        str(getattr(project, "duration", "")) if getattr(project, "duration", None) else None,
    ]
    fallback_text = " ".join(f for f in fallback_fields if f and str(f).strip())

    if not (rfp_text.strip() or fallback_text.strip()):
        return {
            "overview": {
                "Project Name": "Untitled Project",
                "Domain": "TBD",
                "Complexity": "TBD",
                "Tech Stack": "TBD",
                "Use Cases": "TBD",
                "Compliance": "TBD",
                "Duration": 1,
            },
            "activities": [],
            "resourcing_plan": [],
            "architecture_diagram": None,
        }

    kb_results = _rag_retrieve(rfp_text or fallback_text)
    kb_chunks = []
    stop = False
    for group in kb_results:
        for ch in group["chunks"]:
            chunk_tokens = len(tokenizer.encode(ch["content"]))
            if used_tokens + chunk_tokens > max_total_tokens:
                stop = True
                break
            kb_chunks.append(ch["content"])
            used_tokens += chunk_tokens
        if stop:
            break

    logger.info(
        f"Final RFP tokens: {len(rfp_tokens)}, KB tokens: {used_tokens - len(rfp_tokens)}, Total: {used_tokens}/{max_total_tokens}"
    )

    # ---------- Load questions.json (if exists) and build Q&A context ----------
    questions_context = None
    try:
        q_blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        if await azure_blob.blob_exists(q_blob_name):
            q_bytes = await azure_blob.download_bytes(q_blob_name)
            q_json = json.loads(q_bytes.decode("utf-8"))

            q_lines = []
            for category in q_json.get("questions", []):
                cat_name = category.get("category", "General")
                q_lines.append(f"### {cat_name}")
                for item in category.get("items", []):
                    q = item.get("question", "").strip()
                    a = item.get("user_understanding", "").strip() or "(unanswered)"
                    comment = item.get("comment", "").strip()
                    line = f"Q: {q}\nA: {a}"
                    if comment:
                        line += f"\nComment: {comment}"
                    q_lines.append(line)

            questions_context = "\n".join(q_lines)
            logger.info(f"✅ Loaded {len(q_lines)} question lines for project {project.id}")
        else:
            logger.info(f"ℹ️ No questions.json found for project {project.id}, skipping Q&A context.")

    except Exception as e:
        logger.warning(f"⚠️ Could not include questions.json context: {e}")
        questions_context = None

    


    # ---------- Build + query ----------
    prompt = _build_scope_prompt(rfp_text, kb_chunks, project, model_name=model_name, questions_context=questions_context)

    try:
        # Step 1: Generate scope via Azure OpenAI
        resp = await anyio.to_thread.run_sync(
            lambda: client.chat.completions.create(
                model=deployment,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
            )
        )
        raw = _extract_json(resp.choices[0].message.content.strip())
        cleaned_scope = await clean_scope(db, raw, project=project)
        # Update project fields from generated overview (just like finalize_scope)
        overview = cleaned_scope.get("overview", {})
        if overview:
            project.name = overview.get("Project Name") or project.name
            project.domain = overview.get("Domain") or project.domain
            project.complexity = overview.get("Complexity") or project.complexity
            project.tech_stack = overview.get("Tech Stack") or project.tech_stack
            project.use_cases = overview.get("Use Cases") or project.use_cases
            project.compliance = overview.get("Compliance") or project.compliance
            project.duration = str(overview.get("Duration") or project.duration)

            try:
                await db.commit()
                await db.refresh(project)
                logger.info(f" Project metadata updated from generated scope for project {project.id}")
            except Exception as e:
                logger.warning(f" Failed to update project metadata: {e}")


        # Step 2: Generate + store architecture diagram
        try:
            blob_base_path = f"{PROJECTS_BASE}/{getattr(project, 'id', 'unknown')}"
            db_file, arch_blob = await generate_architecture(
                db, project, rfp_text, kb_chunks, blob_base_path
            )
            cleaned_scope["architecture_diagram"] = arch_blob or None
        except Exception as e:
            logger.warning(f"Architecture diagram generation failed: {e}")
            cleaned_scope["architecture_diagram"] = None

        # Step 3: Auto-save finalized_scope.json in Azure Blob + DB
        try:
            from sqlalchemy import select
            result = await db.execute(
                select(models.ProjectFile).filter(
                    models.ProjectFile.project_id == project.id,
                    models.ProjectFile.file_name == "finalized_scope.json",
                )
            )
            old_file = result.scalars().first()
            if old_file:
                logger.info(f"♻️ Overwriting existing finalized_scope.json for project {project.id}")
            else:
                old_file = models.ProjectFile(
                    project_id=project.id,
                    file_name="finalized_scope.json",
                )

            blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"

            await azure_blob.upload_bytes(
                json.dumps(cleaned_scope, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
                overwrite=True,  # ✅ crucial: replace atomically in Blob Storage
            )

            old_file.file_path = blob_name
            db.add(old_file)
            await db.commit()
            await db.refresh(old_file)

            logger.info(f"✅ finalized_scope.json overwritten for project {project.id}")

        except Exception as e:
            logger.warning(f" Failed to auto-save finalized_scope.json: {e}")
        return cleaned_scope

    except Exception as e:
        logger.error(f"Azure OpenAI scope generation failed: {e}")
        return {}


async def regenerate_from_instructions(
    db: AsyncSession,
    project: models.Project,
    draft: dict,
    instructions: str
) -> dict:
    """
    Regenerate the project scope from user instructions using a creative AI-guided prompt.
    Enhances activity sequencing, roles, and effort estimates while preserving valid JSON structure.
    """
    logger.info(f"🔁 Regenerating scope for project {project.id} with creative AI response...")

    if not instructions or not instructions.strip():
        cleaned = await clean_scope(db, draft, project=project)
        return {**cleaned, "_finalized": True}

    if client is None or deployment is None:
        logger.warning("Azure OpenAI not configured; skipping creative regeneration")
        cleaned = await clean_scope(db, draft, project=project)
        return {**cleaned, "_finalized": True}

    # ---- Build creative instruction-aware prompt ----
    prompt = f"""
You are an **expert AI project planner and delivery architect**.
You are given:
1️⃣ The current draft project scope (in JSON).
2️⃣ The user’s specific change instructions.

You must:
- Understand the intent behind the instructions (they may be natural language).
- Modify the draft scope accordingly — you may add, remove, or reorder activities, adjust timelines,
  change resources, or update descriptions to reflect the requested changes.
- You can creatively infer missing dependencies or sequencing logic.
- Preserve all schema keys exactly as before (overview, activities, resourcing_plan).
- All dates must remain valid ISO `yyyy-mm-dd`.
- Keep project duration under 12 months and maximize parallel execution.
- Use realistic IT roles (Backend Developer, Data Engineer, QA Analyst, etc.).
- Always return **only valid JSON**.

If the user requests improvements (e.g., "optimize resourcing", "simplify timeline", "add data validation phase"),
you must reflect those explicitly in the output.

User Instructions:
{instructions}

Current Draft Scope (JSON):
{json.dumps(draft, indent=2, ensure_ascii=False)}

Return ONLY valid JSON with updated overview, activities, and resourcing_plan.
"""

    # ---- Query Azure OpenAI creatively ----
    try:
        resp = await anyio.to_thread.run_sync(
            lambda: client.chat.completions.create(
                model=deployment,
                messages=[
                    {"role": "system", "content": "You are a creative yet precise project scoping expert."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.6,  # a bit higher for creative diversity
            )
        )

        raw_text = resp.choices[0].message.content.strip()
        updated_scope = _extract_json(raw_text)
        cleaned = await clean_scope(db, updated_scope, project=project)

    except Exception as e:
        logger.error(f"⚠️ Creative regeneration failed: {e}")
        cleaned = await clean_scope(db, draft, project=project)

    # ---- Update project metadata from overview ----
    overview = cleaned.get("overview", {})
    if overview:
        project.name = overview.get("Project Name") or project.name
        project.domain = overview.get("Domain") or project.domain
        project.complexity = overview.get("Complexity") or project.complexity
        project.tech_stack = overview.get("Tech Stack") or project.tech_stack
        project.use_cases = overview.get("Use Cases") or project.use_cases
        project.compliance = overview.get("Compliance") or project.compliance
        project.duration = str(overview.get("Duration") or project.duration)
        await db.commit()
        await db.refresh(project)
        logger.info(f"✅ Project metadata synced for project {project.id}")

    # ---- Overwrite finalized_scope.json in Blob ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "finalized_scope.json",
        )
    )
    old_file = result.scalars().first() or models.ProjectFile(
        project_id=project.id, file_name="finalized_scope.json"
    )

    blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(cleaned, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )
    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f"✅ Creative finalized_scope.json regenerated for project {project.id}")
    return {**cleaned, "_finalized": True}




async def finalize_scope(
    db: AsyncSession,
    project_id: str,
    scope_data: dict
) -> tuple[models.ProjectFile, dict]:
    """
    Clean and finalize scope JSON, update project metadata, and store finalized_scope.json in blob.
    """

    logger.info(f" Finalizing scope (engine) for project {project_id}...")

    # ---- Load project with company ----
    result = await db.execute(
        select(models.Project)
        .options(selectinload(models.Project.company))
        .filter(models.Project.id == project_id)
    )
    project = result.scalars().first()

    # ---- Clean the draft using project context ----
    cleaned = await clean_scope(db, scope_data, project=project)
    overview = cleaned.get("overview", {})


    # ---- Update project metadata ----
    result = await db.execute(
        select(models.Project)
        .options(selectinload(models.Project.files))
        .filter(models.Project.id == project_id)
    )
    db_project = result.scalars().first()

    if db_project and overview:
        db_project.name = overview.get("Project Name") or db_project.name
        db_project.domain = overview.get("Domain") or db_project.domain
        db_project.complexity = overview.get("Complexity") or db_project.complexity
        db_project.tech_stack = overview.get("Tech Stack") or db_project.tech_stack
        db_project.use_cases = overview.get("Use Cases") or db_project.use_cases
        db_project.compliance = overview.get("Compliance") or db_project.compliance
        db_project.duration = str(overview.get("Duration") or db_project.duration)
        await db.commit()
        await db.refresh(db_project)

    # ---- Remove old finalized scope if exists ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project_id,
            models.ProjectFile.file_name == "finalized_scope.json"
        )
    )
    old_file = result.scalars().first()
    if old_file:
        logger.info(f"♻️ Overwriting existing finalized_scope.json for project {project_id}")
    else:
        old_file = models.ProjectFile(
            project_id=project_id,
            file_name="finalized_scope.json",
        )

    blob_name = f"{PROJECTS_BASE}/{project_id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(cleaned, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )

    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f"✅ Finalized scope overwritten for project {project_id}")
    return old_file, {**cleaned, "_finalized": True}
